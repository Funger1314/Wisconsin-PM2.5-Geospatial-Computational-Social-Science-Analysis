"""Build markdown, DOCX/PDF, and slide artifacts."""

from __future__ import annotations

from pathlib import Path

import pandas as pd
from docx import Document
from docx.shared import Inches
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches as PptInches
from pptx.util import Pt
from reportlab.lib import colors
from reportlab.lib.pagesizes import landscape, letter
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import inch
from reportlab.platypus import Image, PageBreak, Paragraph, SimpleDocTemplate, Spacer

from .paths import PM25_FIELD, REPORTS_DIR, SLIDES_DIR


def write_project_summary(summary_markdown: str) -> Path:
    """Write the short project summary markdown file."""
    path = REPORTS_DIR / "Wisconsin_PM25_Project_Summary.md"
    path.write_text(summary_markdown, encoding="utf-8")
    return path


def build_writing_sample_docx(
    county_df: pd.DataFrame,
    summary_lines: list[str],
    figure_paths: list[Path],
) -> Path:
    """Create the DOCX writing sample."""
    document = Document()
    document.add_heading("Geographic Patterns in County-Level PM2.5 Monitoring Data Across Wisconsin, 2024", level=0)
    sections = {
        "Research question": summary_lines[0],
        "Data and methods": summary_lines[1],
        "Main findings": summary_lines[2],
        "Spatial-statistics caveat": summary_lines[3],
        "Limitations": summary_lines[4],
        "Policy relevance and future work": summary_lines[5],
    }
    for heading, text in sections.items():
        document.add_heading(heading, level=1)
        document.add_paragraph(text)
    document.add_paragraph(
        f"The monitored-county range in this rebuild runs from {county_df[PM25_FIELD].min():.3f} to {county_df[PM25_FIELD].max():.3f} ug/m3."
    )
    for figure_path in figure_paths[:2]:
        if figure_path.exists():
            document.add_picture(str(figure_path), width=Inches(5.6))
    path = REPORTS_DIR / "Wisconsin_PM25_Writing_Sample_2to3_pages.docx"
    document.save(path)
    return path


def build_writing_sample_pdf(summary_lines: list[str], figure_paths: list[Path]) -> Path:
    """Create the PDF writing sample directly for layout reliability."""
    path = REPORTS_DIR / "Wisconsin_PM25_Writing_Sample_2to3_pages.pdf"
    doc = SimpleDocTemplate(str(path), pagesize=letter, rightMargin=0.65 * inch, leftMargin=0.65 * inch, topMargin=0.65 * inch, bottomMargin=0.65 * inch)
    styles = getSampleStyleSheet()
    title_style = styles["Title"]
    heading_style = styles["Heading2"]
    body_style = ParagraphStyle(
        "Body",
        parent=styles["BodyText"],
        fontSize=10,
        leading=13,
        spaceAfter=8,
    )
    story = [
        Paragraph("Geographic Patterns in County-Level PM2.5 Monitoring Data Across Wisconsin, 2024", title_style),
    ]
    headings = [
        "Research question",
        "Data and methods",
        "Main findings",
        "Spatial-statistics caveat",
        "Limitations",
        "Policy relevance and future work",
    ]
    for heading, text in zip(headings, summary_lines):
        story.append(Paragraph(heading, heading_style))
        story.append(Paragraph(text, body_style))
    for figure_path in figure_paths[:2]:
        if figure_path.exists():
            story.append(Spacer(1, 0.12 * inch))
            story.append(Image(str(figure_path), width=5.8 * inch, height=3.4 * inch))
    doc.build(story)
    return path


def _add_title_band(slide, title: str) -> None:
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, PptInches(13.33), PptInches(0.7))
    band.fill.solid()
    band.fill.fore_color.rgb = RGBColor(152, 28, 36)
    band.line.color.rgb = RGBColor(152, 28, 36)
    text_frame = band.text_frame
    paragraph = text_frame.paragraphs[0]
    paragraph.text = title
    paragraph.font.size = Pt(24)
    paragraph.font.bold = True
    paragraph.font.color.rgb = RGBColor(255, 255, 255)


def _add_bullets(slide, lines: list[str], left: float = 0.6, top: float = 1.0, width: float = 5.4, height: float = 5.8) -> None:
    textbox = slide.shapes.add_textbox(PptInches(left), PptInches(top), PptInches(width), PptInches(height))
    frame = textbox.text_frame
    frame.word_wrap = True
    for idx, line in enumerate(lines):
        paragraph = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        paragraph.text = line
        paragraph.level = 0
        paragraph.font.size = Pt(18)
        paragraph.space_after = Pt(10)


def _add_table_lines(
    slide,
    lines: list[str],
    left: float = 0.6,
    top: float = 1.0,
    width: float = 5.4,
    height: float = 2.4,
) -> None:
    textbox = slide.shapes.add_textbox(PptInches(left), PptInches(top), PptInches(width), PptInches(height))
    frame = textbox.text_frame
    frame.word_wrap = False
    for idx, line in enumerate(lines):
        paragraph = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        paragraph.text = line
        paragraph.font.name = "Menlo"
        paragraph.font.size = Pt(11.5)
        paragraph.space_after = Pt(3)


def _add_takeaway(slide, text: str, left: float = 0.6, top: float = 5.3, width: float = 5.4, height: float = 1.3) -> None:
    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, PptInches(left), PptInches(top), PptInches(width), PptInches(height))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(248, 239, 221)
    box.line.color.rgb = RGBColor(152, 28, 36)
    paragraph = box.text_frame.paragraphs[0]
    paragraph.text = text
    paragraph.font.size = Pt(15)
    paragraph.font.bold = True


def build_slide_deck(
    slide_content: list[dict[str, object]],
    figure_map: dict[str, Path],
) -> Path:
    """Create the PowerPoint deck."""
    presentation = Presentation()
    presentation.slide_width = PptInches(13.33)
    presentation.slide_height = PptInches(7.5)
    blank_layout = presentation.slide_layouts[6]
    for spec in slide_content:
        slide = presentation.slides.add_slide(blank_layout)
        _add_title_band(slide, str(spec["title"]))
        if spec.get("bullets"):
            _add_bullets(slide, [str(item) for item in spec["bullets"]])
        if spec.get("table_lines"):
            _add_table_lines(slide, [str(item) for item in spec["table_lines"]])
        if spec.get("takeaway"):
            _add_takeaway(slide, str(spec["takeaway"]))
        figure_key = spec.get("figure")
        if figure_key and figure_key in figure_map and figure_map[figure_key].exists():
            slide.shapes.add_picture(str(figure_map[figure_key]), PptInches(6.2), PptInches(1.1), width=PptInches(6.5))
        if spec.get("footer"):
            footer = slide.shapes.add_textbox(PptInches(0.5), PptInches(7.05), PptInches(12.2), PptInches(0.3))
            paragraph = footer.text_frame.paragraphs[0]
            paragraph.text = str(spec["footer"])
            paragraph.alignment = PP_ALIGN.RIGHT
            paragraph.font.size = Pt(9)
            paragraph.font.color.rgb = RGBColor(90, 90, 90)
    path = SLIDES_DIR / "Wisconsin_PM25_Spatial_Analysis.pptx"
    presentation.save(path)
    return path


def build_slide_pdf(slide_content: list[dict[str, object]]) -> Path:
    """Create a PDF version of the slide content without requiring PowerPoint export tools."""
    path = SLIDES_DIR / "Wisconsin_PM25_Spatial_Analysis.pdf"
    doc = SimpleDocTemplate(str(path), pagesize=landscape(letter), rightMargin=0.45 * inch, leftMargin=0.45 * inch, topMargin=0.4 * inch, bottomMargin=0.4 * inch)
    styles = getSampleStyleSheet()
    title_style = ParagraphStyle("SlideTitle", parent=styles["Heading1"], textColor=colors.HexColor("#981c24"), fontSize=22, leading=26, spaceAfter=12)
    bullet_style = ParagraphStyle("SlideBullet", parent=styles["BodyText"], fontSize=13, leading=17, leftIndent=14, bulletIndent=0, spaceAfter=8)
    table_style = ParagraphStyle("SlideTable", parent=styles["Code"], fontSize=10.5, leading=13, spaceAfter=4)
    takeaway_style = ParagraphStyle("SlideTakeaway", parent=styles["BodyText"], fontSize=14, leading=18, textColor=colors.HexColor("#981c24"), spaceBefore=8, spaceAfter=8)
    story = []
    for index, spec in enumerate(slide_content):
        story.append(Paragraph(str(spec["title"]), title_style))
        for bullet in spec["bullets"]:
            story.append(Paragraph(str(bullet), bullet_style, bulletText="•"))
        for line in spec.get("table_lines", []):
            story.append(Paragraph(str(line).replace(" ", "&nbsp;"), table_style))
        if spec.get("takeaway"):
            story.append(Paragraph(str(spec["takeaway"]), takeaway_style))
        if index < len(slide_content) - 1:
            story.append(PageBreak())
    doc.build(story)
    return path
